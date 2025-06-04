from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from src.reformat import remove_white_background, auto_crop
import shutil
import os
import streamlit as st

# Output configuration
output_file = "logos_presentation.pptx"
backup_path = "logo_backup"


def clear_folder(cache_folder):
    """
    Deletes all files and folders in the specified directory.

    Args:
        cache_folder (str): Path to the folder to clear.
    """
    for filename in os.listdir(cache_folder):
        file_path = os.path.join(cache_folder, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                shutil.rmtree(file_path)
        except Exception as e:
            st.error(f"Failed to delete {file_path}. Reason: {e}")


def configure_ppt_settings(logo_positions: tuple, slide_size: tuple):
    """
    Calculates layout settings for placing logos in a PowerPoint slide.

    Args:
        logo_positions (tuple): (num_cols, num_rows) - Number of logo columns and rows.
        slide_size (tuple): (width_in_inches, height_in_inches) - Dimensions of the slide.

    Returns:
        dict: Dictionary containing configuration for logo placement and sizing.
    """
    num_cols, num_rows = logo_positions
    user_width, user_height = slide_size
    slide_width, slide_height = Inches(user_width), Inches(user_height)

    # Convert height to pixels (96 DPI) and calculate max height per logo
    logo_height = int(user_height * 96 / num_rows / 2)
    slide_pixel_width = user_width * 96
    max_logo_width = int(slide_pixel_width / num_cols)

    # Center each column horizontally across the slide
    col_spacing = slide_width / (num_cols - 1)
    column_centers = []
    current_x = Inches(0.1)
    for _ in range(num_cols):
        column_center = current_x + col_spacing / 2
        column_centers.append(column_center)
        current_x += col_spacing

    # Vertical spacing between rows
    row_spacing = slide_height / (num_rows - 1)

    return {
        "num_cols": num_cols,
        "num_rows": num_rows,
        "slide_width": slide_width,
        "slide_height": slide_height,
        "logo_height": logo_height,
        "max_logo_width": max_logo_width,
        "column_centers": column_centers,
        "row_spacing": row_spacing,
    }


def load_and_process_logos(
    folder,
    num_cols,
    num_rows,
    slide_width,
    slide_height,
    logo_height,
    max_logo_width,
    column_centers,
    row_spacing,
):
    """
    Loads, cleans, resizes, and saves logos for slide layout.

    Args:
        folder (str): Path to folder where processed logos will be saved.
        num_cols (int): Number of logo columns per slide.
        num_rows (int): Number of logo rows per slide.
        slide_width (Inches): Slide width in inches.
        slide_height (Inches): Slide height in inches.
        logo_height (int): Target logo height in pixels.
        max_logo_width (int): Maximum allowed width for logos.
        column_centers (list): List of horizontal positions (inches) for logo placement.
        row_spacing (float): Vertical spacing (inches) between rows.

    Returns:
        list of tuples: (logo path, resized width in px, resized height in px)
    """

    # Filter only supported image file types
    logos = [
        f for f in os.listdir(folder) if f.lower().endswith((".png", ".jpg", ".jpeg"))
    ]

    # Remove any previous content in the folder to avoid mixing outputs
    clear_folder(folder)
    processed_logos = []

    for logo in logos:
        # Load the original logo from backup
        logo_path = os.path.join(backup_path, logo)
        img = Image.open(logo_path)

        # Step 1: Clean up the logo
        img = remove_white_background(
            img
        )  # Removes white pixels (optional transparency)
        img = auto_crop(img)  # Crops surrounding whitespace to focus on the logo

        # Step 2: Determine new size based on fixed height and aspect ratio
        aspect_ratio = img.width / img.height
        new_width = int(logo_height * aspect_ratio)
        new_height = logo_height

        # Step 3: Enforce a max width to prevent overly wide logos
        if new_width > max_logo_width:
            new_width = max_logo_width
            new_height = int(
                new_width / aspect_ratio
            )  # Recalculate height to maintain proportions

        # Step 4: Resize the image with new dimensions
        img = img.resize((new_width, new_height))

        # Step 5: Save cleaned and resized logo to output folder
        save_path = os.path.join(folder, logo)
        img.save(save_path, format="PNG")

        # Record the final image path and its size for layout placement
        processed_logos.append((save_path, new_width, new_height))

    # Optional: sort logos alphabetically by filename for consistent ordering in output
    processed_logos.sort(key=lambda x: os.path.basename(x[0]).lower())

    return processed_logos


def create_powerpoint(
    processed_logos,
    num_cols,
    num_rows,
    slide_width,
    slide_height,
    logo_height,
    max_logo_width,
    column_centers,
    row_spacing,
):
    """
    Creates a PowerPoint presentation with logos arranged in a grid layout.

    Args:
        processed_logos (list): List of (logo_path, width_px, height_px) tuples.
        num_cols (int): Number of columns of logos.
        num_rows (int): Number of rows of logos.
        slide_width (Inches): Width of the slide.
        slide_height (Inches): Height of the slide.
        logo_height (int): Logo height in pixels.
        max_logo_width (int): Maximum logo width in pixels.
        column_centers (list): X-coordinates of column centers.
        row_spacing (float): Vertical spacing between rows.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Add blank slide

    for idx, (logo_path, width_px, height_px) in enumerate(processed_logos):
        col = idx % num_cols
        row = idx // num_cols

        # Prevent overflow beyond the grid
        if row >= num_rows:
            break

        # Convert dimensions from pixels to inches
        width_in = Inches(width_px / 96)
        height_in = Inches(height_px / 96)

        # Calculate placement coordinates
        x = column_centers[col] - (width_in / 2)
        y = (row + 1) * row_spacing - height_in / 2

        slide.shapes.add_picture(logo_path, x, y, width=width_in, height=height_in)

    prs.save(output_file)
